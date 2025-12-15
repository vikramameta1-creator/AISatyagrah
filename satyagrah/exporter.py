from __future__ import annotations
import csv, io, os, zipfile, math, datetime as dt
from pathlib import Path
from typing import Iterable, Dict, Any, List

from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from reportlab.platypus import Table, TableStyle
from reportlab.lib import colors
from pptx import Presentation
from pptx.util import Inches
try:
    from PIL import Image, ImageDraw
except Exception:
    Image = None
try:
    import imageio.v3 as iio
except Exception:
    iio = None

def _ensure_dir(p: Path) -> None:
    p.parent.mkdir(parents=True, exist_ok=True)

def _sample_rows(n: int = 24) -> List[Dict[str, Any]]:
    now = dt.datetime.utcnow().replace(microsecond=0)
    rows = []
    for i in range(n):
        rows.append({"id": i+1, "timestamp": (now-dt.timedelta(minutes=i)).isoformat()+"Z",
                     "value": round(100*math.sin(i/3)+200, 2), "status": "ok" if i%4 else "peak"})
    return rows

def write_csv(path: Path, rows: Iterable[Dict[str, Any]]) -> None:
    _ensure_dir(path)
    rows = list(rows)
    fields = list(rows[0].keys()) if rows else ["id","timestamp","value","status"]
    with path.open("w", newline="", encoding="utf-8") as f:
        w = csv.DictWriter(f, fieldnames=fields); w.writeheader()
        for r in rows: w.writerow(r)

def write_pdf(path: Path, rows: Iterable[Dict[str, Any]], title="AISatyagrah Export") -> None:
    _ensure_dir(path); rows = list(rows)
    c = canvas.Canvas(str(path), pagesize=A4); W,H = A4
    c.setTitle(title); c.setFont("Helvetica-Bold", 18); c.drawString(72, H-72, title)
    c.setFont("Helvetica", 10); c.drawRightString(W-72, H-68, dt.datetime.utcnow().strftime("%Y-%m-%d %H:%M UTC"))
    if rows:
        data = [list(rows[0].keys())] + [list(r.values()) for r in rows]
        tbl = Table(data, colWidths=[(W-144)/len(data[0])]*len(data[0]))
        tbl.setStyle(TableStyle([("BACKGROUND",(0,0),(-1,0),colors.lightgrey),
                                 ("GRID",(0,0),(-1,-1),0.25,colors.grey),
                                 ("FONTSIZE",(0,0),(-1,-1),9), ("ALIGN",(0,0),(-1,-1),"LEFT")]))
        tbl.wrapOn(c,72,H-120); tbl.drawOn(c,72,max(72,H-140-18*(len(data)+1)))
    c.showPage(); c.save()

def write_pptx(path: Path, rows: Iterable[Dict[str, Any]], title="AISatyagrah Export") -> None:
    _ensure_dir(path); rows = list(rows); prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[0]); s.shapes.title.text = title
    s.placeholders[1].text = dt.datetime.utcnow().strftime("%Y-%m-%d %H:%M UTC")
    if rows:
        s = prs.slides.add_slide(prs.slide_layouts[5])
        R,C = len(rows)+1, len(rows[0]); left, top, width, height = Inches(0.5), Inches(1.3), Inches(9), Inches(5)
        tbl = s.shapes.add_table(R,C,left,top,width,height).table
        hdrs = list(rows[0].keys())
        for j,h in enumerate(hdrs): cell = tbl.cell(0,j); cell.text = str(h); cell.text_frame.paragraphs[0].font.bold = True
        for i,r in enumerate(rows,1):
            for j,h in enumerate(hdrs): tbl.cell(i,j).text = str(r.get(h,""))
    prs.save(str(path))

def write_gif(path: Path, rows: Iterable[Dict[str, Any]], size=(640,360)) -> None:
    _ensure_dir(path)
    if Image is None:
        path.write_bytes(b"GIF89a\x01\x00\x01\x00\x80\x00\x00\x00\x00\x00\xFF\xFF\xFF!\xF9\x04\x01\n\x00\x01\x00,\x00\x00\x00\x00\x01\x00\x01\x00\x00\x02\x02D\x01\x00;"); return
    rows = list(rows) or _sample_rows(20); frames=[]
    for i,r in enumerate(rows):
        img = Image.new("RGB", size, (20,24,28)); d = ImageDraw.Draw(img)
        d.text((20,20), f"Frame {i+1}\nVal: {r['value']}", fill=(220,220,220)); frames.append(img)
    frames[0].save(path, save_all=True, append_images=frames[1:], duration=120, loop=0)

def write_mp4(path: Path, rows: Iterable[Dict[str, Any]], size=(640,360)) -> None:
    _ensure_dir(path)
    if iio is None or Image is None: path.write_bytes(b""); return
    rows = list(rows) or _sample_rows(60); frames=[]
    for i,r in enumerate(rows):
        img = Image.new("RGB", size, (12,18,26)); d = ImageDraw.Draw(img)
        d.text((20,20), f"MP4 Frame {i+1}\nVal: {r['value']}", fill=(230,230,230))
        frames.append(iio.asarray(img))
    iio.imwrite(path, frames, fps=15, codec="libx264", quality=6)

def write_zip(path: Path, files: dict[str, Path]) -> None:
    _ensure_dir(path)
    with zipfile.ZipFile(path,"w",compression=zipfile.ZIP_DEFLATED) as zf:
        for name, fp in files.items():
            if Path(fp).exists(): zf.write(fp, arcname=name)

def build_all_exports(root: Path, when: dt.date|None=None) -> dict:
    when = when or dt.date.today()
    outdir = root / "exports" / when.isoformat()
    ts = dt.datetime.utcnow().strftime("%Y-%m-%d_%H%M%S")
    rows = _sample_rows(28)
    pdf  = outdir / f"export_{ts}.pdf"
    pptx = outdir / f"export_{ts}.pptx"
    csv_ = outdir / f"export_{ts}.csv"
    gif  = outdir / f"export_{ts}.gif"
    mp4  = outdir / f"export_{ts}.mp4"
    zip_ = outdir / f"export_{ts}.zip"

    write_pdf(pdf, rows); write_pptx(pptx, rows); write_csv(csv_, rows)
    write_gif(gif, rows); write_mp4(mp4, rows)
    write_zip(zip_, {pdf.name: pdf, pptx.name: pptx, csv_.name: csv_, gif.name: gif, mp4.name: mp4})

    return {"ok": True, "date": when.isoformat(),
            "pdf": str(pdf).replace("\\","/"), "pptx": str(pptx).replace("\\","/"),
            "csv": str(csv_).replace("\\","/"), "gif": str(gif).replace("\\","/"),
            "mp4": str(mp4).replace("\\","/"), "zip": str(zip_).replace("\\","/")}
