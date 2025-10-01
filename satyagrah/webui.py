# -*- coding: utf-8 -*-
import argparse, sys, subprocess, os, secrets, json, time, io, zipfile, mimetypes, csv, tempfile
from pathlib import Path

try:
    from fastapi import FastAPI, Request
    from fastapi.responses import (
        HTMLResponse, RedirectResponse, JSONResponse,
        StreamingResponse, PlainTextResponse, FileResponse
    )
    from fastapi.staticfiles import StaticFiles
except Exception as e:
    raise SystemExit("FastAPI not installed. Run: pip install fastapi uvicorn") from e

# tolerate missing router package
try:
    from .web.auth_router import router as auth_router
except Exception:
    auth_router = None

from .auth.service import user_from_session
from .core.status import get_status
from .core.jobs import start_job
from .peer.jobfmt import make_job_dict, write_job_zip

# optional telegram notify
try:
    from .notify.telegram import send as telegram_send
except Exception:
    telegram_send = None

app = FastAPI(title="AISatyagrah")

# attach /auth routes if available
if auth_router:
    app.include_router(auth_router, prefix="/auth")

root = Path(__file__).resolve().parents[1]   # satyagrah/
proj = root.parent                           # D:\AISatyagrah

PREFS_PATH     = proj / "data" / "web_prefs.json"
PRESETS_PATH   = proj / "data" / "prompt_presets.json"
FACTS_PATH     = proj / "data" / "facts" / "facts.json"
PENDING_DIR    = proj / "jobs" / "pending"
DEFAULT_REGION = os.getenv("SATYAGRAH_DEFAULT_REGION", "india")

# optional static dir
static_dir = root / "static"
if static_dir.exists():
    app.mount("/static", StaticFiles(directory=str(static_dir)), name="static")

def _html(body: str) -> HTMLResponse:
    return HTMLResponse(f"""<!doctype html>
<html><head><meta charset="utf-8"><title>AISatyagrah</title>
<style>
  :root{{
    --bg:#ffffff; --fg:#111; --muted:#6b7280;
    --card:#fff; --border:#e5e7eb; --shadow:0 3px 12px #0001;
    --btn-bg:#fafafa; --btn-bd:#ddd; --btn-fg:#111; --btn-bg-hover:#f0f0f0;
    --overlay:#111; --overlay-fg:#eee;
  }}
  body.hc{{
    --bg:#0b0b0b; --fg:#fff; --muted:#d1d5db;
    --card:#111; --border:#444; --btn-bg:#1f2937; --btn-bd:#9ca3af; --btn-fg:#fff; --btn-bg-hover:#374151;
    --overlay:#000; --overlay-fg:#f5f5f5;
  }}
  body.no-scroll{{ overflow:hidden; }}

  body{{font-family:system-ui,-apple-system,Segoe UI,Roboto,sans-serif;margin:2rem; background:var(--bg); color:var(--fg)}}
  .grid{{display:grid;grid-template-columns:repeat(auto-fit,minmax(220px,1fr));gap:12px;margin:12px 0}}
  .card{{padding:1rem 1.25rem;border:1px solid var(--border);border-radius:14px;box-shadow:var(--shadow); background:var(--card)}}
  .row{{display:flex;gap:.5rem;flex-wrap:wrap;margin:.25rem 0; align-items:center}}
  .btn{{display:inline-block;padding:.5rem .75rem;border-radius:10px;border:1px solid var(--btn-bd);text-decoration:none;background:var(--btn-bg);cursor:pointer;color:var(--btn-fg)}}
  .btn:hover{{background:var(--btn-bg-hover)}}
  input,textarea,select{{font:inherit;color:var(--fg)}}
  input[type=number]{{width:6rem;padding:.35rem .5rem;border-radius:8px;border:1px solid var(--btn-bd);background:var(--card)}}
  textarea{{width:100%;min-height:80px;padding:.5rem;border-radius:10px;border:1px solid var(--btn-bd);background:var(--card)}}
  .chip{{display:inline-block;padding:.25rem .5rem;border:1px solid var(--btn-bd);border-radius:999px;background:var(--btn-bg);cursor:pointer;margin:.15rem .25rem 0 0;color:var(--btn-fg)}}
  .k{{color:var(--muted)}}
  .v{{font-weight:600}}
  .ok{{color:#059669}}
  .bad{{color:#dc2626}}
  .msg{{min-height:1.2em}}
  code{{background:#f6f6f6;padding:2px 6px;border-radius:6px}}
  table{{width:100%;border-collapse:collapse}}
  th,td{{padding:.4rem .25rem;border-bottom:1px solid var(--border);vertical-align:top}}
  .thumb{{width:144px;height:auto;border-radius:8px;border:1px solid var(--border);cursor:pointer}}

  /* Lightbox */
  .lb-bg{{position:fixed;inset:0;background:rgba(0,0,0,.75);display:none;align-items:center;justify-content:center;z-index:9999}}
  .lb{{background:var(--overlay);border-radius:12px;padding:12px;max-width:90vw;max-height:90vh}}
  .lb img{{max-width:85vw;max-height:75vh;display:block;margin:0 auto;border-radius:8px}}
  .lb-bar{{display:flex;gap:.5rem;justify-content:space-between;align-items:center;margin-top:8px;color:var(--overlay-fg);flex-wrap:wrap}}
  .lb a{{color:var(--overlay-fg);text-decoration:underline}}
  .lb .btn{{color:var(--overlay-fg);background:transparent;border-color:#888}}
  .lb .btn:hover{{background:#222}}

  /* Export dropdown menu */
  .menu{{ position:relative; display:inline-block; }}
  .menu-items{{
    position:absolute; right:0; top:100%; margin-top:6px; min-width:220px;
    background:var(--card); border:1px solid var(--border); border-radius:10px;
    box-shadow:0 10px 30px rgba(0,0,0,.2); padding:6px; display:none; z-index:10000;
  }}
  .menu.open .menu-items{{ display:block; }}
  .menu-item{{
    display:block; width:100%; text-align:left; padding:.5rem .75rem; border-radius:8px;
    background:transparent; border:0; color:var(--fg); cursor:pointer;
  }}
  .menu-item:hover{{ background:var(--btn-bg-hover); }}
</style></head><body>{body}</body></html>""")

def _current_user(request: Request):
    token = request.cookies.get("satyagrah_session")
    return user_from_session(token) if token else None

def _has_role(user, *allowed):
    return bool(user and user.get("role") in allowed)

# ---------- prefs helpers ----------
def _prefs_read() -> dict:
    try:
        if PREFS_PATH.exists():
            return json.loads(PREFS_PATH.read_text(encoding="utf-8"))
    except Exception:
        pass
    return {}

def _prefs_write(data: dict):
    PREFS_PATH.parent.mkdir(parents=True, exist_ok=True)
    PREFS_PATH.write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")

def _get_user_prefs(username: str) -> dict:
    allp = _prefs_read()
    p = allp.get(username) or {}
    return {
        "region": p.get("region") or DEFAULT_REGION,
        "autodrop_default": bool(p.get("autodrop_default", False)),
        "default_steps": int(p.get("default_steps", 28)),
        "default_width": int(p.get("default_width", 768)),
        "default_height": int(p.get("default_height", 1024)),
        "default_count": int(p.get("default_count", 1)),
    }

def _set_user_prefs(username: str, newvals: dict):
    allp = _prefs_read()
    cur = allp.get(username) or {}
    allowed = {"region", "autodrop_default", "default_steps", "default_width", "default_height", "default_count"}
    for k,v in newvals.items():
        if k in allowed:
            cur[k] = v
    allp[username] = cur
    _prefs_write(allp)

# ---------- presets / facts helpers ----------
def _load_presets() -> list[str]:
    if PRESETS_PATH.exists():
        try:
            data = json.loads(PRESETS_PATH.read_text(encoding="utf-8"))
            if isinstance(data, list):
                return [str(x) for x in data][:40]
        except Exception:
            pass
    return [
        "low-poly political poster, bold palette, high contrast",
        "newspaper cutout collage, satire headline style",
        "pop-art halftone comic panel, speech bubble",
        "vintage propaganda poster remix, witty tagline",
        "street-art stencil look, gritty texture",
        "retro-futurism chrome, neon rim light",
        "papercraft diorama, tiny crowd, soft shadows",
        "minimal flat illustration, geometric shapes"
    ]

def _prompt_from_facts(region: str) -> str:
    try:
        data = json.loads(FACTS_PATH.read_text(encoding="utf-8"))
        topics = data.get("topics", [])[:4]
        bits = []
        for t in topics:
            title = (t.get("title") or "").strip()
            tags  = (t.get("tags") or [])[:3]
            if title:
                bits.append(f"{title} ({', '.join(tags)})")
        if bits:
            joined = "; ".join(bits)
            return f"satirical poster about: {joined}; region {region}; witty but tasteful; strong composition"
    except Exception:
        pass
    return f"political satire poster; region {region}; bold colors; clean layout"

# ---------- results helpers ----------
def _results_dir() -> Path|None:
    p = os.getenv("SATYAGRAH_PEER_RESULTS", "")
    return Path(p) if p else None

def _result_zip(job_id: str) -> Path|None:
    resdir = _results_dir()
    if not resdir:
        return None
    zp = resdir / f"result_{job_id}.zip"
    return zp if zp.exists() else None

def _zip_images(zp: Path) -> list[str]:
    with zipfile.ZipFile(zp, "r") as zf:
        return [n for n in zf.namelist() if n.lower().endswith((".png",".jpg",".jpeg",".webp",".gif"))]

# ---------- Jobs API ----------
@app.get("/api/jobs")
def api_jobs(limit: int = 50, offset: int = 0):
    outbox = Path(os.getenv("SATYAGRAH_PEER_OUTBOX", str(proj / "jobs" / "outbox")))
    inbox_env   = os.getenv("SATYAGRAH_PEER_INBOX", "")
    results_env = os.getenv("SATYAGRAH_PEER_RESULTS", "")
    inbox   = Path(inbox_env)   if inbox_env   else None
    results = Path(results_env) if results_env else None

    items = []
    now = time.time()

    try:
        if not outbox.exists():
            return JSONResponse({"ok": True, "items": [], "meta": {
                "outbox": str(outbox), "inbox": str(inbox) if inbox else "",
                "results": str(results) if results else "", "limit": limit, "offset": offset, "count": 0
            }})

        limit = max(1, min(200, int(limit)))
        offset = max(0, int(offset))

        all_zips = sorted(outbox.glob("job_*.zip"),
                          key=lambda p: p.stat().st_mtime,
                          reverse=True)
        subset = all_zips[offset:offset+limit]

        for p in subset:
            st = p.stat()
            job_id = p.stem.replace("job_", "")
            rec = {
                "job_id": job_id,
                "mtime": int(st.st_mtime),
                "age_sec": int(now - st.st_mtime),
                "size": st.st_size,
                "out_zip": str(p),
            }
            if inbox:
                rec["inbox_has_copy"] = (inbox / p.name).exists()

            status = "unknown"
            if results:
                rp = results / f"result_{job_id}.zip"
                if rp.exists():
                    rs = rp.stat()
                    rec.update({
                        "result_found": True,
                        "result_zip": str(rp),
                        "result_mtime": int(rs.st_mtime),
                        "result_age_sec": int(now - rs.st_mtime),
                        "result_size": rs.st_size,
                    })
                    status = "done"
                else:
                    rec["result_found"] = False
                    status = "queued"
            rec["status"] = status
            items.append(rec)

        return JSONResponse({"ok": True, "items": items, "meta": {
            "outbox": str(outbox), "inbox": str(inbox) if inbox else "",
            "results": str(results) if results else "", "limit": limit, "offset": offset, "count": len(items)
        }})
    except Exception as e:
        return JSONResponse({"ok": False, "error": str(e)}, status_code=500)

# ---------- results / thumbs / images / downloads ----------
@app.get("/api/results")
def api_results(limit: int = 24):
    resdir = _results_dir()
    if not resdir or not resdir.exists():
        return JSONResponse({"ok": True, "items": []})
    limit = max(1, min(100, int(limit)))
    zips = sorted(resdir.glob("result_*.zip"), key=lambda p: p.stat().st_mtime, reverse=True)[:limit]
    items = []
    for z in zips:
        job_id = z.stem.replace("result_", "")
        st = z.stat()
        try:
            count = len(_zip_images(z))
        except Exception:
            count = 0
        items.append({"job_id": job_id, "zip": str(z), "mtime": int(st.st_mtime), "size": st.st_size, "image_count": count})
    return JSONResponse({"ok": True, "items": items})

@app.get("/api/result_manifest/{job_id}")
def api_result_manifest(job_id: str):
    zp = _result_zip(job_id)
    if not zp:
        return JSONResponse({"ok": False, "error": "not found"}, status_code=404)
    try:
        with zipfile.ZipFile(zp, "r") as zf:
            rows = []
            for n in zf.namelist():
                info = zf.getinfo(n)
                mime = mimetypes.guess_type(n)[0] or "application/octet-stream"
                rows.append({"name": n, "size": info.file_size, "mime": mime})
        return JSONResponse({"ok": True, "items": rows})
    except Exception as e:
        return JSONResponse({"ok": False, "error": str(e)}, status_code=500)

@app.get("/result_thumb/{job_id}")
def result_thumb(job_id: str):
    zp = _result_zip(job_id)
    if not zp: return JSONResponse({"ok": False, "error": "not found"}, status_code=404)
    try:
        with zipfile.ZipFile(zp, "r") as zf:
            names = [n for n in zf.namelist() if n.lower().endswith((".png",".jpg",".jpeg",".webp",".gif"))]
            if not names: return JSONResponse({"ok": False, "error": "no image in zip"}, status_code=404)
            with zf.open(names[0]) as fh:
                data = fh.read()
                ctype = mimetypes.guess_type(names[0])[0] or "image/png"
                return StreamingResponse(io.BytesIO(data), media_type=ctype)
    except Exception as e:
        return JSONResponse({"ok": False, "error": str(e)}, status_code=500)

@app.get("/result_image/{job_id}/{idx}")
def result_image(job_id: str, idx: int = 0):
    zp = _result_zip(job_id)
    if not zp: return JSONResponse({"ok": False, "error": "not found"}, status_code=404)
    try:
        with zipfile.ZipFile(zp, "r") as zf:
            names = [n for n in zf.namelist() if n.lower().endswith((".png",".jpg",".jpeg",".webp",".gif"))]
            if not names: return JSONResponse({"ok": False, "error": "no image in zip"}, status_code=404)
            idx = max(0, min(len(names)-1, int(idx)))
            with zf.open(names[idx]) as fh:
                data = fh.read()
                ctype = mimetypes.guess_type(names[idx])[0] or "image/png"
                return StreamingResponse(io.BytesIO(data), media_type=ctype)
    except Exception as e:
        return JSONResponse({"ok": False, "error": str(e)}, status_code=500)

@app.get("/download_result_image_file/{job_id}/{idx}")
def download_result_image_file(job_id: str, idx: int = 0):
    zp = _result_zip(job_id)
    if not zp: return JSONResponse({"ok": False, "error": "not found"}, status_code=404)
    try:
        with zipfile.ZipFile(zp, "r") as zf:
            names = [n for n in zf.namelist() if n.lower().endswith((".png",".jpg",".jpeg",".webp",".gif"))]
            if not names: return JSONResponse({"ok": False, "error": "no image in zip"}, status_code=404)
            idx = max(0, min(len(names)-1, int(idx)))
            name = Path(names[idx]).name
            with zf.open(names[idx]) as fh:
                data = fh.read()
                ctype = mimetypes.guess_type(name)[0] or "application/octet-stream"
                headers = {"Content-Disposition": f'attachment; filename="{name}"'}
                return StreamingResponse(io.BytesIO(data), media_type=ctype, headers=headers)
    except Exception as e:
        return JSONResponse({"ok": False, "error": str(e)}, status_code=500)

@app.get("/download_result_zip/{job_id}")
def download_result_zip(job_id: str):
    zp = _result_zip(job_id)
    if not zp: return JSONResponse({"ok": False, "error": "not found"}, status_code=404)
    return FileResponse(str(zp), media_type="application/zip", filename=f"result_{job_id}.zip")

@app.get("/download_result_csv/{job_id}")
def download_result_csv(job_id: str):
    zp = _result_zip(job_id)
    if not zp: return JSONResponse({"ok": False, "error": "not found"}, status_code=404)
    try:
        output = io.StringIO()
        w = csv.writer(output)
        w.writerow(["name","size","mime"])
        with zipfile.ZipFile(zp, "r") as zf:
            for n in zf.namelist():
                info = zf.getinfo(n)
                w.writerow([n, info.file_size, mimetypes.guess_type(n)[0] or "application/octet-stream"])
        data = output.getvalue().encode("utf-8")
        headers = {"Content-Disposition": f'attachment; filename="result_{job_id}.csv"'}
        return StreamingResponse(io.BytesIO(data), media_type="text/csv", headers=headers)
    except Exception as e:
        return JSONResponse({"ok": False, "error": str(e)}, status_code=500)

@app.post("/api/extract_result")
async def api_extract_result(request: Request):
    try:
        data = await request.json()
        job_id = str(data.get("job_id") or "").strip()
        open_explorer = bool(data.get("open_explorer", False))
    except Exception:
        return JSONResponse({"ok": False, "error": "invalid json"}, status_code=400)

    zp = _result_zip(job_id)
    if not zp:
        return JSONResponse({"ok": False, "error": "not found"}, status_code=404)

    target = zp.parent / "extracted" / job_id
    try:
        target.mkdir(parents=True, exist_ok=True)
        with zipfile.ZipFile(zp, "r") as zf:
            zf.extractall(target)
        if open_explorer and os.name == "nt":
            try: os.startfile(str(target))  # type: ignore[attr-defined]
            except Exception: pass
        return JSONResponse({"ok": True, "folder": str(target)})
    except Exception as e:
        return JSONResponse({"ok": False, "error": str(e)}, status_code=500)

# ---------- NEW: PDF / PPTX / GIF / MP4 ----------
@app.get("/download_contact_sheet_pdf/{job_id}")
def download_contact_sheet_pdf(job_id: str, cols: int = 3, rows: int = 3, paper: str = "a4", dpi: int = 150, margin: int = 36):
    """Build a contact sheet PDF (grid of images per page)."""
    zp = _result_zip(job_id)
    if not zp: return JSONResponse({"ok": False, "error": "not found"}, status_code=404)

    try:
        from PIL import Image, ImageOps
    except Exception:
        return JSONResponse({"ok": False, "error": "Pillow not installed. pip install pillow"}, status_code=500)

    paper = paper.lower()
    if paper == "letter":
        W_in, H_in = 8.5, 11.0
    else:
        W_in, H_in = 8.27, 11.69

    page_w = int(W_in * dpi)
    page_h = int(H_in * dpi)
    cols = max(1, int(cols)); rows = max(1, int(rows))
    cell_w = (page_w - 2*margin) // cols
    cell_h = (page_h - 2*margin) // rows

    try:
        with zipfile.ZipFile(zp, "r") as zf:
            names = [n for n in zf.namelist() if n.lower().endswith((".png",".jpg",".jpeg",".webp",".gif"))]
            if not names:
                return JSONResponse({"ok": False, "error": "no images"}, status_code=404)

            pages = []
            page = Image.new("RGB", (page_w, page_h), "white")
            i = 0
            for n in names:
                with zf.open(n) as fh:
                    im = Image.open(io.BytesIO(fh.read()))
                    if im.mode in ("RGBA", "LA"):
                        bg = Image.new("RGB", im.size, "white")
                        bg.paste(im, mask=im.getchannel("A") if "A" in im.getbands() else None)
                        im = bg
                    im = ImageOps.exif_transpose(im)
                    im.thumbnail((cell_w-8, cell_h-8))
                    x = margin + (i % cols) * cell_w + (cell_w - im.width)//2
                    y = margin + (i // cols % rows) * cell_h + (cell_h - im.height)//2
                    page.paste(im, (x, y))
                    i += 1
                    if (i % (cols*rows)) == 0:
                        pages.append(page)
                        page = Image.new("RGB", (page_w, page_h), "white")
            if i % (cols*rows) != 0:
                pages.append(page)

            buff = io.BytesIO()
            pages[0].save(buff, format="PDF", save_all=True, append_images=pages[1:])
            buff.seek(0)
            headers = {"Content-Disposition": f'attachment; filename="result_{job_id}_contact_sheet.pdf"'}
            return StreamingResponse(buff, media_type="application/pdf", headers=headers)
    except Exception as e:
        return JSONResponse({"ok": False, "error": str(e)}, status_code=500)

@app.get("/download_result_pptx/{job_id}")
def download_result_pptx(job_id: str):
    """Build a PPTX (one slide per image). Requires python-pptx."""
    zp = _result_zip(job_id)
    if not zp: return JSONResponse({"ok": False, "error": "not found"}, status_code=404)

    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception:
        return JSONResponse({"ok": False, "error": "python-pptx not installed. pip install python-pptx"}, status_code=500)

    try:
        prs = Presentation()
        blank = prs.slide_layouts[6]
        with zipfile.ZipFile(zp, "r") as zf:
            names = [n for n in zf.namelist() if n.lower().endswith((".png",".jpg",".jpeg",".webp",".gif"))]
            if not names:
                return JSONResponse({"ok": False, "error": "no images"}, status_code=404)
            for n in names:
                slide = prs.slides.add_slide(blank)
                data = zf.read(n)
                pic = slide.shapes.add_picture(io.BytesIO(data), Inches(0), Inches(0))
                slide_w, slide_h = prs.slide_width, prs.slide_height
                ratio = min(slide_w / pic.width, slide_h / pic.height)
                pic.width = int(pic.width * ratio)
                pic.height = int(pic.height * ratio)
                pic.left = int((slide_w - pic.width) / 2)
                pic.top  = int((slide_h - pic.height) / 2)
        buff = io.BytesIO()
        prs.save(buff); buff.seek(0)
        headers = {"Content-Disposition": f'attachment; filename="result_{job_id}.pptx"'}
        return StreamingResponse(buff, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", headers=headers)
    except Exception as e:
        return JSONResponse({"ok": False, "error": str(e)}, status_code=500)

@app.get("/download_result_gif/{job_id}")
def download_result_gif(job_id: str, fps: float = 2.0, loop: int = 0):
    """Animated GIF slideshow using Pillow only."""
    zp = _result_zip(job_id)
    if not zp: return JSONResponse({"ok": False, "error": "not found"}, status_code=404)
    try:
        from PIL import Image, ImageOps, ImagePalette  # noqa
    except Exception:
        return JSONResponse({"ok": False, "error": "Pillow not installed. pip install pillow"}, status_code=500)

    try:
        with zipfile.ZipFile(zp, "r") as zf:
            names = [n for n in zf.namelist() if n.lower().endswith((".png",".jpg",".jpeg",".webp",".gif"))]
            if not names:
                return JSONResponse({"ok": False, "error": "no images"}, status_code=404)
            frames = []
            for n in names:
                im = Image.open(io.BytesIO(zf.read(n)))
                im = ImageOps.exif_transpose(im)
                if im.mode in ("RGBA", "LA"):
                    bg = Image.new("RGB", im.size, "white"); bg.paste(im, mask=im.split()[-1])
                    im = bg
                frames.append(im.convert("P", palette=Image.ADAPTIVE))
            dur = int(max(1, 1000.0 / max(0.1, float(fps))))
            buff = io.BytesIO()
            frames[0].save(buff, format="GIF", save_all=True, append_images=frames[1:], duration=dur, loop=max(0,int(loop)))
            buff.seek(0)
            headers = {"Content-Disposition": f'attachment; filename="result_{job_id}.gif"'}
            return StreamingResponse(buff, media_type="image/gif", headers=headers)
    except Exception as e:
        return JSONResponse({"ok": False, "error": str(e)}, status_code=500)

@app.get("/download_result_mp4/{job_id}")
def download_result_mp4(job_id: str, fps: float = 2.0):
    """MP4 slideshow (requires moviepy + imageio-ffmpeg)."""
    zp = _result_zip(job_id)
    if not zp: return JSONResponse({"ok": False, "error": "not found"}, status_code=404)
    try:
        import moviepy.editor as mpy
    except Exception:
        return JSONResponse({"ok": False, "error": "MP4 export requires moviepy and imageio-ffmpeg. pip install moviepy imageio-ffmpeg"}, status_code=500)

    try:
        with zipfile.ZipFile(zp, "r") as zf:
            names = [n for n in zf.namelist() if n.lower().endswith((".png",".jpg",".jpeg",".webp",".gif"))]
            if not names:
                return JSONResponse({"ok": False, "error": "no images"}, status_code=404)

            with tempfile.TemporaryDirectory() as td:
                paths = []
                for i,n in enumerate(names):
                    p = Path(td) / f"frame_{i:04d}.png"
                    p.write_bytes(zf.read(n))
                    paths.append(str(p))
                clip = mpy.ImageSequenceClip(paths, fps=max(0.1,float(fps)))
                outp = Path(td) / f"result_{job_id}.mp4"
                clip.write_videofile(str(outp), codec="libx264", audio=False, verbose=False, logger=None)
                return FileResponse(str(outp), media_type="video/mp4", filename=outp.name)
    except Exception as e:
        return JSONResponse({"ok": False, "error": str(e)}, status_code=500)

# ---------- presets / facts endpoints ----------
@app.get("/api/presets")
def api_presets():
    return JSONResponse({"ok": True, "items": _load_presets()})

@app.get("/api/prompt_from_facts")
def api_prompt_from_facts(region: str = DEFAULT_REGION):
    return JSONResponse({"ok": True, "prompt": _prompt_from_facts(region)})

# ---------- prefs API ----------
@app.get("/api/prefs")
def api_prefs(request: Request):
    u = _current_user(request)
    if not u: return JSONResponse({"ok": False, "error": "unauthorized"}, status_code=403)
    return JSONResponse({"ok": True, "prefs": _get_user_prefs(u["username"])})

@app.post("/api/prefs")
async def api_prefs_set(request: Request):
    u = _current_user(request)
    if not u: return JSONResponse({"ok": False, "error": "unauthorized"}, status_code=403)
    try:
        data = await request.json()
    except Exception:
        return JSONResponse({"ok": False, "error": "invalid json"}, status_code=400)
    try:
        _set_user_prefs(u["username"], data)
        return JSONResponse({"ok": True, "prefs": _get_user_prefs(u["username"])})
    except Exception as e:
        return JSONResponse({"ok": False, "error": str(e)}, status_code=500)

# ---------- viewer approvals ----------
def _pending_dir() -> Path:
    PENDING_DIR.mkdir(parents=True, exist_ok=True)
    return PENDING_DIR

@app.get("/api/pending")
def api_pending():
    pend = _pending_dir()
    zips = sorted(pend.glob("job_*.zip"), key=lambda p: p.stat().st_mtime, reverse=True)
    items = [{"job_id": p.stem.replace("job_", ""), "zip": str(p), "mtime": int(p.stat().st_mtime), "size": p.stat().st_size} for p in zips]
    return JSONResponse({"ok": True, "items": items})

@app.post("/api/approve_job")
async def api_approve_job(request: Request):
    u = _current_user(request)
    if not _has_role(u, "admin", "editor"):
        return JSONResponse({"ok": False, "error": "unauthorized"}, status_code=403)
    try:
        data = await request.json()
        job_id = str(data.get("job_id"))
    except Exception:
        return JSONResponse({"ok": False, "error": "invalid json"}, status_code=400)
    src = _pending_dir() / f"job_{job_id}.zip"
    if not src.exists():
        return JSONResponse({"ok": False, "error": "not found"}, status_code=404)
    outbox = Path(os.getenv("SATYAGRAH_PEER_OUTBOX", str(proj / "jobs" / "outbox")))
    outbox.mkdir(parents=True, exist_ok=True)
    dst = outbox / src.name
    dst.write_bytes(src.read_bytes())
    try: src.unlink()
    except Exception: pass
    dropped = False
    inbox_env = os.getenv("SATYAGRAH_PEER_INBOX", "")
    if inbox_env:
        try:
            inbox = Path(inbox_env); inbox.mkdir(parents=True, exist_ok=True)
            (inbox / dst.name).write_bytes(dst.read_bytes())
            dropped = True
        except Exception:
            dropped = False
    return JSONResponse({"ok": True, "out_zip": str(dst), "dropped": dropped})

# ---------- Pages ----------
@app.get("/", response_class=HTMLResponse)
def index(request: Request):
    u = _current_user(request)
    if u:
        body = f"""
        <div class="card">
          <h2>Welcome, {u['username']} 👋</h2>
          <p>Role: <b>{u['role']}</b></p>
          <div class="row">
            <a class="btn" href="/dash">Open Dashboard</a>
            <form method="post" action="/auth/logout"><button class="btn" type="submit">Logout</button></form>
            <a class="btn" href="/auth/me">/auth/me</a>
          </div>
        </div>"""
    else:
        body = """
        <div class="card">
          <h2>Welcome 👋</h2>
          <p>You’re not signed in.</p>
          <div class="row"><a class="btn" href="/auth/login">Login</a></div>
        </div>"""
    return _html(body)

@app.get("/dash", response_class=HTMLResponse)
def dashboard(request: Request):
    u = _current_user(request)
    if not u:
        return RedirectResponse(url="/auth/login")

    s = get_status()
    is_admin_editor = _has_role(u, "admin", "editor")
    is_viewer       = _has_role(u, "viewer")

    # IMPORTANT: f-string to interpolate tile values
    tiles = f"""
    <div class="grid">
      <div class="card">
        <div class="k">SD host</div>
        <div class="v">{s['sd_host']}</div>
        <div class="{ 'ok' if s['sd_reachable'] else 'bad'}">{'reachable' if s['sd_reachable'] else 'down'}</div>
      </div>
      <div class="card">
        <div class="k">Secret</div>
        <div class="v">{'set ✅' if s['secret_set'] else 'missing ❌'}</div>
      </div>
      <div class="card">
        <div class="k">Telegram</div>
        <div class="v">{'configured ✅' if s['telegram_configured'] else 'not configured'}</div>
      </div>
      <div class="card">
        <div class="k">Latest Run</div>
        <div class="v">{s['latest_run_date'] or '-'}</div>
        <div class="k">Exports</div>
        <div class="v">{s['exports_count_for_latest']}</div>
      </div>
    </div>
    """

    prefs_card = ""
    if is_admin_editor:
        prefs_card = """
        <div class="card">
          <h3>Preferences</h3>
          <div class="row">
            <label>Region&nbsp;
              <select id="region">
                <option value="india">India</option>
                <option value="usa">USA</option>
                <option value="uk">UK</option>
                <option value="eu">EU</option>
                <option value="global">Global</option>
              </select>
            </label>
            <label><input id="autodrop_default" type="checkbox"> Auto-drop jobs to agent inbox</label>
          </div>
          <div class="row">
            <label>Default Steps <input id="p_steps" type="number" value="28"></label>
            <label>Default W <input id="p_w" type="number" value="768"></label>
            <label>Default H <input id="p_h" type="number" value="1024"></label>
            <label>Default Count <input id="p_count" type="number" value="1"></label>
          </div>
          <div class="row"><button class="btn" onclick="savePrefs()">Save preferences</button></div>
          <div id="p_msg" class="k msg"></div>
        </div>
        """

    outbox = os.getenv("SATYAGRAH_PEER_OUTBOX", str(proj / "jobs" / "outbox"))
    inbox  = os.getenv("SATYAGRAH_PEER_INBOX", "")
    inbox_hint = f"<div class='k'>Auto-drop inbox: <code>{inbox}</code></div>" if inbox else "<div class='k'>Set <code>SATYAGRAH_PEER_INBOX</code> to auto-drop jobs to an agent.</div>"
    creator_title = "Peer Job — Create &amp; Sign" if is_admin_editor else "Request a Peer Job (needs approval)"
    create_btn = "Make job.zip" if is_admin_editor else "Request job"
    autodrop_btn = "Make & auto-drop to inbox" if is_admin_editor else ""
    disable_auto = "" if is_admin_editor else "style='display:none'"

    actions = f"""
    <div class="card">
      <h3>{creator_title}</h3>
      <div class="k">Outbox: <code>{outbox}</code></div>
      {inbox_hint if is_admin_editor else ""}
      <div class="k">Presets:</div>
      <div id="chips" class="row"></div>
      <div class="row" style="margin:.25rem 0">
        <button class="btn" onclick="useFacts()">Use today’s facts</button>
      </div>
      <textarea id="prompt" placeholder="Describe the image..."></textarea>
      <div class="row">
        <label>Seed <input id="seed" type="number" value="77"></label>
        <label>Steps <input id="steps" type="number" value="28"></label>
        <label>W <input id="w" type="number" value="768"></label>
        <label>H <input id="h" type="number" value="1024"></label>
        <label>Count <input id="count" type="number" value="1" min="1" max="8"></label>
      </div>
      <div class="row">
        <button class="btn" onclick="makeJob()">{create_btn}</button>
        <button class="btn" {disable_auto} onclick="makeJob(true)">{autodrop_btn}</button>
      </div>
      <div id="msg" class="k"></div>
    </div>
    """

    pending_card = ""
    if is_admin_editor:
        pending_card = """
        <div class="card">
          <h3>Pending requests</h3>
          <table><thead><tr><th>Job ID</th><th>Created</th><th>Action</th></tr></thead>
          <tbody id="pendBody"><tr><td colspan="3" class="k">Loading…</td></tr></tbody></table>
        </div>
        <script>
        async function loadPending(){
          try{
            const r = await fetch('/api/pending'); const j = await r.json();
            const tb = document.getElementById('pendBody');
            if(!j.ok){ tb.innerHTML = '<tr><td colspan=3 class=k>Error</td></tr>'; return; }
            if(!j.items.length){ tb.innerHTML = '<tr><td colspan=3 class=k>None</td></tr>'; return; }
            tb.innerHTML = j.items.map(x=>{
              const t=new Date(x.mtime*1000).toLocaleString();
              return `<tr><td><code>${x.job_id}</code></td><td>${t}</td>
                <td><button class="btn" onclick="approve('${x.job_id}')">Approve</button></td></tr>`;
            }).join('');
          }catch(e){ document.getElementById('pendBody').innerHTML='<tr><td colspan=3 class=k>Error</td></tr>'; }
          setTimeout(loadPending,5000);
        }
        async function approve(id){
          const r = await fetch('/api/approve_job',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({job_id:id})});
          await r.json(); loadPending();
        }
        loadPending();
        </script>
        """

    gallery_and_more = """
    <div class="card">
      <h3>Recent peer jobs</h3>
      <div class="k">Shows job_*.zip in outbox and result_*.zip (if SATYAGRAH_PEER_RESULTS is set).</div>
      <table id="jobsTbl">
        <thead><tr style="text-align:left;border-bottom:1px solid var(--border)">
          <th>Job ID</th><th>Created</th><th>Status</th><th>Outbox</th><th>Result</th></tr></thead>
        <tbody id="jobsBody"><tr><td colspan="5" class="k">Loading…</td></tr></tbody>
      </table>
    </div>

    <div class="card">
      <h3>Result gallery</h3>
      <div id="gal" class="row"></div>
    </div>

    <div class="lb-bg" id="lbBg">
      <div class="lb">
        <img id="lbImg" src="" alt="">
        <div class="lb-bar">
          <div class="row" style="gap:.5rem;flex:1 1 auto">
            <div id="lbCaption" class="k">-</div>
          </div>
          <div class="row" style="gap:.5rem">
            <button class="btn" onclick="lbPrev()">◀ Prev</button>
            <button class="btn" onclick="lbNext()">Next ▶</button>

            <div class="menu">
              <button class="btn" onclick="toggleExportMenu()" aria-haspopup="true" aria-expanded="false" id="exportBtn">Export ▾</button>
              <div class="menu-items" id="exportMenu" role="menu">
                <button class="menu-item" role="menuitem" onclick="exportAs('zip')">ZIP (all)</button>
                <button class="menu-item" role="menuitem" onclick="exportAs('image')">Image (current)</button>
                <hr style="border:none;border-top:1px solid var(--border);margin:.25rem 0">
                <button class="menu-item" role="menuitem" onclick="exportAs('csv')">CSV (manifest)</button>
                <button class="menu-item" role="menuitem" onclick="exportAs('pdf')">PDF (contact sheet)</button>
                <button class="menu-item" role="menuitem" onclick="exportAs('pptx')">PPTX (slides)</button>
                <button class="menu-item" role="menuitem" onclick="exportAs('gif')">GIF (slideshow)</button>
                <button class="menu-item" role="menuitem" onclick="exportAs('mp4')">MP4 (slideshow)</button>
                <hr style="border:none;border-top:1px solid var(--border);margin:.25rem 0">
                <button class="menu-item" role="menuitem" onclick="exportAs('extract')">Extract to folder</button>
              </div>
            </div>

            <button class="btn" onclick="lbClose()">Close</button>
          </div>
        </div>
      </div>
    </div>

    <div class="card">
      <h3>PeerAgent — Service mode</h3>
      <div class="k">Download a PowerShell script to install/remove a scheduled task that starts the PeerAgent at logon.</div>
      <div class="row">
        <label>PeerAgent.exe path <input id="paPath" style="width:28rem" value="C:\\PeerAgent\\PeerAgent.exe"></label>
      </div>
      <div class="row">
        <button class="btn" onclick="dlService('create')">Download “Install at logon” .ps1</button>
        <button class="btn" onclick="dlService('remove')">Download “Remove task” .ps1</button>
      </div>
    </div>

    <script>
let PREFS = { region: 'india', autodrop_default: false, default_steps: 28,
              default_width: 768, default_height: 1024, default_count: 1 };
const LB = {job:'', idx:0, count:1};

/* THEME */
function applyTheme(hc){
  document.body.classList.toggle('hc', !!hc);
  try{ localStorage.setItem('satyagrah_hc',''+(hc?1:0)); }catch(e){}
}
(function initTheme(){
  try{
    const hc = localStorage.getItem('satyagrah_hc')==='1';
    document.body.classList.toggle('hc', hc);
    const t = document.getElementById('hcToggle'); if(t) t.checked = hc;
  }catch(e){}
})();

async function loadPrefs(){
  try{
    const r = await fetch('/api/prefs'); const j = await r.json();
    if(j.ok){ PREFS = j.prefs || PREFS; }
  }catch(e){}
  const R = document.getElementById('region'); if (R) R.value = PREFS.region || 'india';
  const AD = document.getElementById('autodrop_default'); if (AD) AD.checked = !!PREFS.autodrop_default;
  const S  = document.getElementById('p_steps');  if (S)  S.value  = PREFS.default_steps  || 28;
  const W  = document.getElementById('p_w');      if (W)  W.value  = PREFS.default_width  || 768;
  const H  = document.getElementById('p_h');      if (H)  H.value  = PREFS.default_height || 1024;
  const C  = document.getElementById('p_count');  if (C)  C.value  = PREFS.default_count  || 1;
}

async function savePrefs(){
  const body = {
    region: (document.getElementById('region')||{value:'india'}).value,
    autodrop_default: !!(document.getElementById('autodrop_default')||{checked:false}).checked,
    default_steps: +(document.getElementById('p_steps')||{value:28}).value,
    default_width: +(document.getElementById('p_w')||{value:768}).value,
    default_height:+(document.getElementById('p_h')||{value:1024}).value,
    default_count: +(document.getElementById('p_count')||{value:1}).value
  };
  try{
    const r = await fetch('/api/prefs',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify(body)});
    const j = await r.json();
    document.getElementById('p_msg').textContent = j.ok ? 'Saved.' : 'Error';
    if(j.ok){ await loadPrefs(); }
  }catch(e){ document.getElementById('p_msg').textContent = 'Error'; }
}

async function loadPresets(){
  try{
    const r=await fetch('/api/presets'); const j=await r.json(); if(!j.ok) return;
    const chips=document.getElementById('chips');
    chips.innerHTML = j.items.map(t=>`<span class="chip" onclick="addPreset(\`${t.replace(/`/g,'\\`')}\`)">${t}</span>`).join('');
  }catch(e){}
}
function addPreset(t){ const p=document.getElementById('prompt'); p.value=(p.value? p.value+' ': '')+t; }

async function useFacts(){
  const reg=(document.getElementById('region')||{value:'india'}).value;
  const r=await fetch('/api/prompt_from_facts?region='+encodeURIComponent(reg)); const j=await r.json();
  if(j.ok){ const p=document.getElementById('prompt'); p.value=j.prompt; }
}

async function makeJob(autodropButton=false){
  const body={
    prompt:document.getElementById('prompt').value,
    seed:+document.getElementById('seed').value,
    steps:+document.getElementById('steps').value,
    width:+document.getElementById('w').value,
    height:+document.getElementById('h').value,
    count:+document.getElementById('count').value,
    region:(document.getElementById('region')||{value:'india'}).value,
    autodrop:autodropButton || (!!(document.getElementById('autodrop_default')||{checked:false}).checked)
  };
  const r=await fetch('/api/make_job',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify(body)});
  const j=await r.json();
  document.getElementById('msg').textContent = j.ok
    ? (j.pending ? ('Queued for approval: '+j.pending_zip)
                 : ('Wrote: '+j.out_zip+(j.dropped?' — and dropped to inbox ✅':'')))
    : ('Error: '+j.error);
  loadJobs(); loadResults();
}

async function loadJobs(){
  try{
    const r=await fetch('/api/jobs?limit=50'); const j=await r.json(); const tb=document.getElementById('jobsBody');
    if(!j.ok){ tb.innerHTML='<tr><td colspan=5 class=k>Error</td></tr>'; return; }
    if(!j.items.length){ tb.innerHTML='<tr><td colspan=5 class=k>No jobs yet.</td></tr>'; return; }
    tb.innerHTML=j.items.map(x=>{
      const created=new Date(x.mtime*1000).toLocaleString();
      const out=(x.out_zip||'').replaceAll('\\\\','/');
      const res=x.result_found?`✅ (${(x.result_zip||'').replaceAll('\\\\','/')})`:'—';
      return `<tr><td><code>${x.job_id}</code></td><td>${created}</td><td>${x.status||'-'}</td><td><code>${out}</code></td><td>${res}</td></tr>`;
    }).join('');
  }catch(e){ document.getElementById('jobsBody').innerHTML='<tr><td colspan=5 class=k>Error</td></tr>'; }
  setTimeout(loadJobs,5000);
}

async function loadResults(){
  try{
    const r=await fetch('/api/results?limit=24'); const j=await r.json(); const g=document.getElementById('gal');
    if(!j.ok||!j.items.length){ g.innerHTML='<span class=k>No results yet.</span>'; return; }
    g.innerHTML=j.items.map(x=>`<div style="text-align:center">
      <img class="thumb" src="/result_thumb/${x.job_id}" alt="${x.job_id}"
           onclick="lbOpen('${x.job_id}',0,${x.image_count||1})">
      <div class="k">${x.job_id} (${x.image_count||1})</div>
    </div>`).join('');
  }catch(e){}
  setTimeout(loadResults,6000);
}

function lbOpen(job,idx,count){
  LB.job=job; LB.idx=idx||0; LB.count=Math.max(1, count||1);
  document.getElementById('lbImg').src='/result_image/'+job+'/'+LB.idx;
  document.getElementById('lbCaption').textContent=job+'  ['+(LB.idx+1)+'/'+LB.count+']';
  document.getElementById('lbBg').style.display='flex';
  document.body.classList.add('no-scroll');  // lock background scroll
}
function lbClose(){
  document.getElementById('lbBg').style.display='none';
  document.body.classList.remove('no-scroll');
}
function lbNext(){ LB.idx=(LB.idx+1)%LB.count; lbOpen(LB.job, LB.idx, LB.count); }
function lbPrev(){ LB.idx=(LB.idx-1+LB.count)%LB.count; lbOpen(LB.job, LB.idx, LB.count); }

/* Export menu logic */
function toggleExportMenu(){
  const m = document.getElementById('exportMenu').parentElement;
  const open = m.classList.toggle('open');
  document.getElementById('exportBtn').setAttribute('aria-expanded', open?'true':'false');
}
document.addEventListener('click', (e)=>{
  const menu = document.getElementById('exportMenu')?.parentElement;
  if(!menu) return;
  if(!menu.contains(e.target)) menu.classList.remove('open');
});
async function exportAs(kind){
  const menu = document.getElementById('exportMenu')?.parentElement;
  if(menu) menu.classList.remove('open');
  const job = LB.job || ''; if(!job) return;

  if(kind==='zip'){ window.open('/download_result_zip/'+job, '_blank'); return; }
  if(kind==='image'){ window.open('/download_result_image_file/'+job+'/'+(LB.idx||0), '_blank'); return; }
  if(kind==='csv'){ window.open('/download_result_csv/'+job, '_blank'); return; }
  if(kind==='pdf'){ window.open('/download_contact_sheet_pdf/'+job+'?cols=3&rows=3&paper=a4&dpi=150', '_blank'); return; }
  if(kind==='pptx'){ window.open('/download_result_pptx/'+job, '_blank'); return; }
  if(kind==='gif'){ window.open('/download_result_gif/'+job+'?fps=2', '_blank'); return; }
  if(kind==='mp4'){ window.open('/download_result_mp4/'+job+'?fps=2', '_blank'); return; }
  if(kind==='extract'){
    try{
      const r = await fetch('/api/extract_result',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({job_id: job, open_explorer: true})});
      await r.json();
    }catch(e){}
    return;
  }
}

/* Service-mode script */
async function dlService(action){
  const exe = (document.getElementById('paPath')||{value:'C:\\\\PeerAgent\\\\PeerAgent.exe'}).value;
  const url = '/api/peeragent_service_ps1?action='+encodeURIComponent(action)+'&exe_path='+encodeURIComponent(exe);
  const txt = await (await fetch(url)).text();
  const a = document.createElement('a');
  a.href = URL.createObjectURL(new Blob([txt],{type:'text/plain'}));
  a.download = (action==='create'?'PeerAgent_install.ps1':'PeerAgent_remove.ps1');
  a.click();
  setTimeout(()=>URL.revokeObjectURL(a.href), 1000);
}

/* Init */
(async()=>{ await loadPrefs(); await loadPresets(); loadJobs(); loadResults(); })();
    </script>
    """

    body = f"""
    <div class="card">
      <h2>Dashboard</h2>
      <div class="row" style="width:100%;justify-content:space-between">
        <div class="row">
          <a class="btn" href="/">Home</a>
          <form method="post" action="/auth/logout"><button class="btn" type="submit">Logout</button></form>
        </div>
        <label class="k" style="display:flex;align-items:center;gap:.5rem">
          <span>High contrast</span>
          <input type="checkbox" id="hcToggle" onchange="applyTheme(this.checked)">
        </label>
      </div>
      <p>Signed in as <b>{u['username']}</b> · role <b>{u['role']}</b></p>
    </div>
    {tiles}
    {prefs_card if is_admin_editor else ""}
    {actions}
    {pending_card if is_admin_editor else ""}
    {gallery_and_more}
    {"<div class='card'><p>Read-only mode: your requests go to Pending for admin approval.</p></div>" if is_viewer else ""}
    """
    return _html(body)

# ---------- actions & make_job ----------
def _popen(args: list[str]):
    return subprocess.Popen(args, cwd=str(proj))

@app.post("/actions/captions")
def action_captions(request: Request):
    u = _current_user(request)
    if not _has_role(u, "admin", "editor"):
        return RedirectResponse(url="/auth/login")
    region = _get_user_prefs(u["username"]).get("region", DEFAULT_REGION)
    cmd = [sys.executable, "-m", "satyagrah.captions.after_batch",
           "--date", "latest", "--top", "3", "--region", region]
    start_job(cmd, name=f"captions-{region}")
    return RedirectResponse(url="/dash")

@app.post("/actions/batch")
def action_batch(request: Request):
    u = _current_user(request)
    if not _has_role(u, "admin", "editor"):
        return RedirectResponse(url="/auth/login")
    cmd = [sys.executable, "-m", "satyagrah", "batch", "--top", "3", "--date", "latest", "--seed", "12345", "--package"]
    start_job(cmd, name="batch")
    return RedirectResponse(url="/dash")

@app.get("/api/status")
def api_status():
    return JSONResponse(get_status())

@app.post("/api/make_job")
async def api_make_job(request: Request):
    u = _current_user(request)
    is_admin_editor = _has_role(u, "admin", "editor")
    is_viewer       = _has_role(u, "viewer")
    if not (is_admin_editor or is_viewer):
        return JSONResponse({"ok": False, "error": "unauthorized"}, status_code=403)

    try:
        data = await request.json()
    except Exception:
        return JSONResponse({"ok": False, "error": "invalid json"}, status_code=400)

    prompt = (data.get("prompt") or "").strip()
    if not prompt:
        return JSONResponse({"ok": False, "error": "prompt required"}, status_code=400)

    seed   = int(data.get("seed", 77))
    steps  = int(data.get("steps", 28))
    width  = int(data.get("width", 768))
    height = int(data.get("height", 1024))
    count  = int(data.get("count", 1))
    region = (data.get("region") or _get_user_prefs(u["username"])["region"])
    autodrop = bool(data.get("autodrop", False)) if is_admin_editor else False

    job_id = f"demo_{secrets.token_hex(4)}"
    tasks = [{"type":"txt2img","prompt":prompt,"seed":seed,"steps":steps,"width":width,"height":height,"count":count}]
    job = make_job_dict(job_id, requester=u["username"], tasks=tasks, ttl_hours=24)
    job.setdefault("meta", {})["region"] = region

    if is_viewer:
        pend = _pending_dir() / f"job_{job_id}.zip"
        try:
            write_job_zip(job, pend)
        except Exception as e:
            return JSONResponse({"ok": False, "error": str(e)})
        return JSONResponse({"ok": True, "pending": True, "pending_zip": str(pend)})

    outbox = Path(os.getenv("SATYAGRAH_PEER_OUTBOX", str(proj / "jobs" / "outbox")))
    outbox.mkdir(parents=True, exist_ok=True)
    out_zip = outbox / f"job_{job_id}.zip"
    try:
        write_job_zip(job, out_zip)
    except Exception as e:
        return JSONResponse({"ok": False, "error": str(e)})

    dropped = False
    if autodrop:
        try:
            inbox = Path(os.getenv("SATYAGRAH_PEER_INBOX", ""))
            if str(inbox):
                inbox.mkdir(parents=True, exist_ok=True)
                (inbox / out_zip.name).write_bytes(out_zip.read_bytes())
                dropped = True
        except Exception:
            dropped = False

    try:
        if telegram_send:
            telegram_send(f"AISatyagrah: new peer job {job_id} [{region}] ({width}x{height}, steps={steps}, count={count})")
    except Exception:
        pass

    return JSONResponse({"ok": True, "job_id": job_id, "out_zip": str(out_zip), "dropped": dropped})

# ---------- PeerAgent “service mode” script ----------
@app.get("/api/peeragent_service_ps1")
def api_peeragent_service_ps1(action: str = "create", exe_path: str = r"C:\PeerAgent\PeerAgent.exe"):
    task_name = "AISatyagrahPeerAgent"
    if action not in {"create","remove"}:
        return JSONResponse({"ok": False, "error": "action must be create|remove"}, status_code=400)
    if action == "create":
        script = f"""$ErrorActionPreference='Stop'
$task = New-ScheduledTaskAction -Execute "{exe_path}"
$trig = New-ScheduledTaskTrigger -AtLogOn
Register-ScheduledTask -TaskName "{task_name}" -Action $task -Trigger $trig -Description "AISatyagrah PeerAgent" -Force
Write-Host "Installed task '{task_name}' to start at logon."
"""
    else:
        script = f"""$ErrorActionPreference='Stop'
Unregister-ScheduledTask -TaskName "{task_name}" -Confirm:$false
Write-Host "Removed task '{task_name}'."
"""
    return PlainTextResponse(script)

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--host", default="127.0.0.1")
    parser.add_argument("--port", type=int, default=8000)
    args = parser.parse_args()
    import uvicorn
    uvicorn.run("satyagrah.webui:app", host=args.host, port=args.port, reload=False)

if __name__ == "__main__":
    main()
