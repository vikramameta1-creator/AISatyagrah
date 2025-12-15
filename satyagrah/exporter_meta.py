# -*- coding: utf-8 -*-
"""
Satyagraph meta exporters (PDF + PPTX)

- Reads topic rows via facts_satire.build_topic_rows(date)
- Enriches with LoRA jokes stored in data/runs/<date>/satire.json (field: lora_joke)
- Writes:
    - PDF: compact overview (title, summary, one-liner, LoRA joke)
    - PPTX: one slide per topic, with both one-liner and LoRA joke when present

Usage examples:

  from pathlib import Path
  from satyagrah.exporter_meta import write_pdf_for_topics, write_pptx_for_topics

  run_date = "2025-09-18"
  root = Path("data/runs") / run_date

  write_pdf_for_topics(str(root / "satyagraph_meta.pdf"), run_date)
  write_pptx_for_topics(str(root / "satyagraph_meta.pptx"), run_date)
"""

from __future__ import annotations

import os
import json
from pathlib import Path
from typing import Any, Dict, List, Optional

from .analysis.facts_satire import build_topic_rows


def _root_default() -> Path:
    d = Path(r"D:\AISatyagrah")
    return d if d.exists() else Path.cwd()


ROOT: Path = Path(os.environ.get("AISATYAGRAH_ROOT") or _root_default()).resolve()
DATA: Path = ROOT / "data"
RUNS: Path = DATA / "runs"


# ------------------------ helpers ------------------------


def _wrap(text: str, width: int = 90) -> List[str]:
    """
    Simple word-wrap helper for PDF text.
    """
    text = (text or "").strip()
    if not text:
        return []
    words = text.split()
    lines: List[str] = []
    cur: List[str] = []
    for w in words:
        cur_len = sum(len(x) for x in cur) + max(0, len(cur) - 1)
        if cur_len + len(w) > width and cur:
            lines.append(" ".join(cur))
            cur = [w]
        else:
            cur.append(w)
    if cur:
        lines.append(" ".join(cur))
    return lines


def _load_lora_joke_map(run_date: str) -> Dict[str, str]:
    """
    Reads data/runs/<date>/satire.json and returns {topic_id: lora_joke}.
    """
    satire_path = RUNS / run_date / "satire.json"
    if not satire_path.exists():
        return {}
    try:
        with satire_path.open("r", encoding="utf-8") as f:
            satire = json.load(f)
    except Exception:
        return {}

    if not isinstance(satire, dict):
        return {}

    out: Dict[str, str] = {}
    for tid, entry in satire.items():
        if not isinstance(entry, dict):
            continue
        lj = entry.get("lora_joke")
        if isinstance(lj, str) and lj.strip():
            out[str(tid)] = lj.strip()
    return out


def _attach_lora_jokes(run_date: str, rows: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    """
    Attach 'lora_joke' from satire.json into each row if available.
    """
    lmap = _load_lora_joke_map(run_date)
    if not lmap:
        return rows

    out: List[Dict[str, Any]] = []
    for r in rows:
        tid = str(r.get("topic_id") or "")
        lj = lmap.get(tid)
        if lj:
            r2 = dict(r)
            r2["lora_joke"] = lj
            out.append(r2)
        else:
            out.append(r)
    return out


def _ensure_parent(path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)


# ------------------------ PDF exporter ------------------------


def write_pdf_for_topics(out_path: str, run_date: str) -> str:
    """
    Create a simple PDF meta report for a given run date.

    Includes:
      - Title
      - Neutral summary
      - One-liner
      - LoRA joke (if present)
    """
    out = Path(out_path)
    _ensure_parent(out)

    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.pdfgen import canvas
    except Exception as e:
        out.write_text(
            f"Satyagraph meta PDF exporter requires reportlab.\nError: {e}\n",
            encoding="utf-8",
        )
        print(f"[exporter_meta] reportlab not available, wrote placeholder note to {out}")
        return str(out)

    try:
        rows = build_topic_rows(run_date)
        rows = _attach_lora_jokes(run_date, rows)
    except Exception as e:
        out.write_text(
            f"Error building topic rows for {run_date}: {e}\n", encoding="utf-8"
        )
        print(f"[exporter_meta] build_topic_rows failed for {run_date}: {e}")
        return str(out)

    c = canvas.Canvas(str(out), pagesize=A4)
    width, height = A4

    c.setTitle(f"Satyagraph meta – {run_date}")

    def new_page(title_suffix: str = ""):
        c.showPage()
        y = height - 40
        c.setFont("Helvetica-Bold", 16)
        title = f"Satyagraph – {run_date}"
        if title_suffix:
            title += f" ({title_suffix})"
        c.drawString(40, y, title)
        y -= 30
        c.setFont("Helvetica", 9)
        return y

    # First page header
    y = height - 40
    c.setFont("Helvetica-Bold", 16)
    c.drawString(40, y, f"Satyagraph – {run_date}")
    y -= 30
    c.setFont("Helvetica", 9)

    page_num = 1

    for idx, row in enumerate(rows, start=1):
        title = (row.get("title") or "").strip() or f"Topic {row.get('topic_id')}"
        summary = (row.get("summary") or "").strip()
        one = (row.get("one_liner") or "").strip()
        lora = (row.get("lora_joke") or "").strip()

        # estimate needed lines; if low, start a new page
        est_lines = 3  # for title + spacing
        if summary:
            est_lines += len(_wrap("Summary: " + summary))
        if one:
            est_lines += len(_wrap("One-liner: " + one))
        if lora:
            est_lines += len(_wrap("LoRA joke: " + lora))
        if y < 60 + est_lines * 11:
            page_num += 1
            y = new_page(f"page {page_num}")

        # Title
        c.setFont("Helvetica-Bold", 11)
        c.drawString(40, y, f"{idx}. {title}")
        y -= 14
        c.setFont("Helvetica", 9)

        # Text fields
        lines: List[str] = []
        if summary:
            lines.append("Summary: " + summary)
        if one:
            lines.append("One-liner: " + one)
        if lora:
            lines.append("LoRA joke: " + lora)

        for text in lines:
            for ln in _wrap(text, width=100):
                c.drawString(50, y, ln)
                y -= 11
        y -= 8  # extra spacing between topics

    c.save()
    print(f"[exporter_meta] wrote PDF to {out}")
    return str(out)


# ------------------------ PPTX exporter ------------------------


def write_pptx_for_topics(out_path: str, run_date: str) -> str:
    """
    Create a PPTX where each topic is one slide:

    Title       -> slide title
    Summary     -> main body text
    One-liner   -> bullet
    LoRA joke   -> bullet (if present)
    """
    out = Path(out_path)
    _ensure_parent(out)

    try:
        from pptx import Presentation
        from pptx.util import Pt
    except Exception as e:
        out.write_text(
            f"Satyagraph meta PPTX exporter requires python-pptx (and Pillow).\nError: {e}\n",
            encoding="utf-8",
        )
        print(f"[exporter_meta] python-pptx not available, wrote placeholder note to {out}")
        return str(out)

    try:
        rows = build_topic_rows(run_date)
        rows = _attach_lora_jokes(run_date, rows)
    except Exception as e:
        out.write_text(
            f"Error building topic rows for {run_date}: {e}\n", encoding="utf-8"
        )
        print(f"[exporter_meta] build_topic_rows failed for {run_date}: {e}")
        return str(out)

    prs = Presentation()

    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Satyagraph"
    subtitle = slide.placeholders[1]
    subtitle.text = f"Meta overview – {run_date}"

    # Content slides
    content_layout = prs.slide_layouts[1]

    for row in rows:
        title = (row.get("title") or "").strip() or f"Topic {row.get('topic_id')}"
        summary = (row.get("summary") or "").strip()
        one = (row.get("one_liner") or "").strip()
        lora = (row.get("lora_joke") or "").strip()

        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = title

        body = slide.placeholders[1].text_frame
        body.clear()

        if summary:
            p = body.add_paragraph()
            p.text = summary
            p.level = 0
            p.font.size = Pt(18)

        if one:
            p = body.add_paragraph()
            p.text = "One-liner: " + one
            p.level = 1
            p.font.size = Pt(14)

        if lora:
            p = body.add_paragraph()
            p.text = "LoRA joke: " + lora
            p.level = 1
            p.font.size = Pt(14)

        if not summary and not one and not lora:
            p = body.add_paragraph()
            p.text = "(no content for this topic)"
            p.level = 0
            p.font.size = Pt(16)

    prs.save(str(out))
    print(f"[exporter_meta] wrote PPTX to {out}")
    return str(out)
