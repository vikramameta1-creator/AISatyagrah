from pathlib import Path
from typing import List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from ..storage import captions as capstore

def _find_images(date: str, root: Path) -> List[Path]:
    hits: List[Path] = []
    for base in [root / "exports" / date, root / "data" / "runs" / date, root / "data" / "runs" / date / "art"]:
        if base.exists():
            hits += sorted(p for p in base.rglob("*") if p.suffix.lower() in (".png", ".jpg", ".jpeg"))
    return hits[:120]

def run(*, date: str, exports_root: Path, files: Optional[List[str]] = None, **_) -> Path:
    root = exports_root.parent
    outdir = exports_root / date
    outdir.mkdir(parents=True, exist_ok=True)
    path = outdir / "export.pptx"

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(10), Inches(1.2))
    tf = title.text_frame; tf.text = f"AISatyagrah — {date}"; tf.paragraphs[0].font.size = Pt(40)

    imgs = [Path(f) for f in (files or []) if Path(f).exists()] or _find_images(date, root)
    caps = capstore.load(root, date)

    for p in imgs:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        left, top = Inches(0.5), Inches(0.5)
        width, height = Inches(12.33), Inches(6.0)
        slide.shapes.add_picture(str(p), left, top, width=width, height=height)

        rel = p.relative_to(root).as_posix()
        meta = caps.get(rel, {})
        text = ((meta.get("caption") or "") + " " + (meta.get("hashtags") or "")).strip()
        if text:
            tb = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(12.33), Inches(0.7))
            t = tb.text_frame; t.clear()
            p0 = t.paragraphs[0]; p0.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
            run = p0.add_run()
            run.text = text
            run.font.size = Pt(16)

    prs.save(str(path))
    return path
