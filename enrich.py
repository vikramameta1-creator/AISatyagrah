from __future__ import annotations
from pathlib import Path
from datetime import datetime
from typing import Dict, List
import os, subprocess, shutil

# --- helpers you already drafted ---
def make_pdf(file_paths: List[str], out_path: str, title: str="AISatyagrah Export") -> str:
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.units import cm
        from reportlab.pdfgen import canvas
    except Exception:
        Path(out_path).with_suffix(".txt").write_text(
            "Install reportlab for real PDFs:  pip install reportlab\n", encoding="utf-8")
        return out_path
    c = canvas.Canvas(out_path, pagesize=A4)
    w,h = A4
    y = h - 2*cm
    c.setFont("Helvetica-Bold",16); c.drawString(2*cm,y,title); y-=1.2*cm
    c.setFont("Helvetica",10); c.drawString(2*cm,y,f"Generated: {datetime.now().isoformat(timespec='seconds')}"); y-=1.0*cm
    for p in file_paths:
        if y < 3*cm:
            c.showPage(); y = h - 2*cm
        name = os.path.basename(p)
        c.setFont("Helvetica-Bold",12); c.drawString(2*cm,y,name); y-=0.5*cm
        c.setFont("Helvetica",10); c.drawString(2*cm,y,p); y-=0.7*cm
    c.save(); return out_path

def make_pptx(image_paths: List[str], out_path: str, title: str="AISatyagrah Export") -> str:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from PIL import Image
    except Exception:
        Path(out_path).with_suffix(".txt").write_text(
            "Install python-pptx & Pillow:  pip install python-pptx Pillow\n", encoding="utf-8")
        return out_path
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = tx.text_frame; tf.text = title; tf.paragraphs[0].font.size = Pt(28)
    x,y,col = 0.5,1.1,0
    for p in image_paths:
        try:
            Image.open(p).close()
        except Exception:
            continue
        slide.shapes.add_picture(p, Inches(x), Inches(y), width=Inches(3.2))
        col += 1; x += 3.4
        if col==3:
            col=0; x=0.5; y+=2.6
            if y>6.5:
                slide = prs.slides.add_slide(prs.slide_layouts[5]); y=0.5
    prs.save(out_path); return out_path

# --- main entry required by the API ---
def export_all(root: Path, date: str, kind: str) -> Dict[str,str]:
    """
    root: exports root (Path), e.g. D:\\AISatyagrah\\exports
    date: 'YYYY-MM-DD'
    kind: 'all' | 'pdf' | 'pptx' | 'csv' | 'gif' | 'mp4' | 'zip'
    Returns: mapping of artifact-name -> absolute path (strings)
    """
    outdir = Path(root) / date
    outdir.mkdir(parents=True, exist_ok=True)
    stamp = datetime.now().strftime("%Y-%m-%d_%H%M%S")

    # Gather inputs (simple heuristic: files under 'outdir')
    all_files = sorted([str(p) for p in outdir.rglob("*") if p.is_file()])
    images = [p for p in all_files if os.path.splitext(p)[1].lower() in {".jpg",".jpeg",".png",".webp"}]

    artifacts: Dict[str,str] = {}

    # CSV listing (always handy)
    csv_path = outdir / f"export_{date}_{stamp}.csv"
    try:
        csv_path.write_text("path,size\n" + "\n".join(f"{p},{os.path.getsize(p)}" for p in all_files), encoding="utf-8")
    except Exception:
        pass
    artifacts["csv"] = str(csv_path)

    # PDF
    if kind in ("all","pdf"):
        pdf_path = outdir / f"export_{date}_{stamp}.pdf"
        artifacts["pdf"] = make_pdf(all_files, str(pdf_path), title=f"AISatyagrah Export — {date}")

    # PPTX
    if kind in ("all","pptx"):
        pptx_path = outdir / f"export_{date}_{stamp}.pptx"
        artifacts["pptx"] = make_pptx(images, str(pptx_path), title=f"AISatyagrah Export — {date}")

    # Lightweight GIF/MP4 placeholders (upgrade later to real ffmpeg calls)
    if kind in ("all","gif"):
        gif_path = outdir / f"export_{date}_{stamp}.gif"
        Path(gif_path).write_bytes(b"GIF89a")   # tiny valid header; replace with real render later
        artifacts["gif"] = str(gif_path)

    if kind in ("all","mp4"):
        mp4_path = outdir / f"export_{date}_{stamp}.mp4"
        mp4_path.write_text("Install ffmpeg to render MP4s\n", encoding="utf-8")
        artifacts["mp4"] = str(mp4_path)

    # Zip everything we produced
    if kind in ("all","zip"):
        zip_base = outdir / f"export_{date}_{stamp}"
        # include all artifacts we made so far
        temp_dir = outdir / f"__pkg_{stamp}"
        temp_dir.mkdir(exist_ok=True)
        for k,v in artifacts.items():
            try:
                src = Path(v)
                if src.exists():
                    shutil.copy2(src, temp_dir / src.name)
            except Exception:
                pass
        zip_path = shutil.make_archive(str(zip_base), "zip", root_dir=temp_dir)
        artifacts["zip"] = zip_path
        try:
            shutil.rmtree(temp_dir)
        except Exception:
            pass

    return artifacts
